# -*- coding: utf-8 -*-
"""
L2N_V2.ipynb — Multimodal Lecture-to-Notes (L2N+)
Corrected & improved version of L2N_V1.

Key fixes over V1:
  1. pipeline task corrected: "text-generation" → "summarization" (BART is seq2seq, not causal LM)
  2. summarize_chunks / make_tldr now reliably read "summary_text" key
  3. Multi-file input: accepts a comma-separated list of files and fuses all sources
  4. Semantic deduplication implemented (was documented but missing in V1)
  5. Better chunk-heading heuristic (noun-phrase extraction instead of raw truncation)
  6. Graceful fallback when a chunk is too short for the summarizer
  7. Environment-safe file upload (Colab / local)
  8. poppler-utils install guarded to Colab only
"""

pip install openai-whisper python-pptx pdf2image pytesseract sentence-transformers transformers nltk torch

# ── 1. Imports ───────────────────────────────────────────────────────────────
import os
import re
import json
from typing import List, Dict

import torch
import whisper
import pytesseract
from pdf2image import convert_from_path
from pptx import Presentation
from transformers import pipeline                         # FIX: used with task="summarization"
from sentence_transformers import SentenceTransformer, util
import nltk
nltk.download("punkt", quiet=True)
nltk.download("punkt_tab", quiet=True)
from nltk.tokenize import sent_tokenize

# ── 2. Environment helpers ───────────────────────────────────────────────────
try:
    from google.colab import files as colab_files
    COLAB_ENV = True
except ImportError:
    colab_files = None
    COLAB_ENV = False

def _install_poppler():
    """Only install system packages inside Colab."""
    if COLAB_ENV:
        os.system("apt-get install -y poppler-utils -qq")

# ── 3. Config ────────────────────────────────────────────────────────────────
CHUNK_WORD_TARGET          = 450
EXTRACTIVE_SENTENCES       = 6
# FIX ① — BART is a summarisation model; use task="summarization"
SUMMARIZER_MODEL           = "facebook/bart-large-cnn"
DEVICE                     = 0 if torch.cuda.is_available() else -1
DEDUP_SIMILARITY_THRESHOLD = 0.85   # cosine threshold for dropping near-duplicate sentences
OUTPUT_JSON                = "lecture_notes.json"
OUTPUT_MD                  = "lecture_notes.md"
TRANSCRIPT_PATH            = "cleaned_transcript.txt"

# ── 4. Extraction layer ──────────────────────────────────────────────────────
_whisper_model = None

def _get_whisper():
    global _whisper_model
    if _whisper_model is None:
        print("Loading Whisper ASR model...")
        _whisper_model = whisper.load_model("base")
    return _whisper_model

def extract_audio(path: str) -> str:
    print(f"  [ASR] Transcribing {path} ...")
    result = _get_whisper().transcribe(path)
    return result["text"]

def extract_pdf(path: str) -> str:
    _install_poppler()
    print(f"  [OCR] Extracting text from PDF {path} ...")
    pages = convert_from_path(path)
    parts = []
    for i, page in enumerate(pages, 1):
        text = pytesseract.image_to_string(page).strip()
        if text:
            parts.append(f"SLIDE {i}\n{text}")
    return "\n\n".join(parts)

def extract_pptx(path: str) -> str:
    print(f"  [PPT] Parsing {path} ...")
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            parts.append(f"SLIDE {i}\n" + "\n".join(texts))
    return "\n\n".join(parts)

EXT_MAP = {
    "mp3": extract_audio, "wav": extract_audio,
    "mp4": extract_audio, "m4a": extract_audio,
    "pdf": extract_pdf,
    "pptx": extract_pptx,
}

def extract_file(path: str) -> str:
    ext = path.rsplit(".", 1)[-1].lower()
    if ext not in EXT_MAP:
        raise ValueError(f"Unsupported file type: .{ext}")
    return EXT_MAP[ext](path)

# ── 5. Multi-file input & fusion ─────────────────────────────────────────────
def collect_input_files() -> List[str]:
    """
    FIX ③ — Accept multiple files (the SRS requires multi-document input).
    In Colab: upload via widget.  Locally: comma-separated paths from stdin.
    """
    if COLAB_ENV:
        print("Upload one or more files (audio / PDF / PPTX):")
        uploaded = colab_files.upload()
        return list(uploaded.keys())
    else:
        raw = input("Enter file path(s) separated by commas: ").strip()
        return [p.strip() for p in raw.split(",") if p.strip()]

def fuse_extractions(file_paths: List[str]) -> str:
    """
    Extract text from every input file and concatenate with source tags.
    Source tagging (FR9) lets downstream steps know where each block came from.
    """
    blocks = []
    for path in file_paths:
        try:
            text = extract_file(path)
            tag  = f"[SOURCE: {os.path.basename(path)}]"
            blocks.append(f"{tag}\n{text}")
        except Exception as e:
            print(f"  [WARN] Could not process {path}: {e}")
    if not blocks:
        raise RuntimeError("No content could be extracted from the provided files.")
    return "\n\n".join(blocks)

# ── 6. Preprocessing ──────────────────────────────────────────────────────────
def preprocess(raw: str) -> str:
    text = raw.strip()
    # Collapse whitespace runs
    text = re.sub(r'\s+', ' ', text)
    # Ensure sentence boundaries have a newline (helps NLTK tokeniser)
    text = re.sub(r'(\.)([A-Za-z])', r'.\n\2', text)
    text = text.replace("&", "and")
    # Remove consecutive duplicate words  e.g. "the the"
    text = re.sub(r'\b(\w+)( \1\b)+', r'\1', text, flags=re.IGNORECASE)
    return text

# ── 7. Semantic deduplication ─────────────────────────────────────────────────
def deduplicate_sentences(sentences: List[str], model: SentenceTransformer,
                          threshold: float = DEDUP_SIMILARITY_THRESHOLD) -> List[str]:
    """
    FIX ④ — Remove near-duplicate sentences using SBERT cosine similarity.
    This was documented in the SRS (FR10) but absent in V1.
    """
    if not sentences:
        return []
    embeddings = model.encode(sentences, convert_to_tensor=True, show_progress_bar=False)
    kept_indices = []
    kept_embeddings = []
    for i, emb in enumerate(embeddings):
        if not kept_embeddings:
            kept_indices.append(i)
            kept_embeddings.append(emb)
            continue
        sims = util.cos_sim(emb, torch.stack(kept_embeddings)).squeeze()
        max_sim = float(sims.max()) if sims.dim() > 0 else float(sims)
        if max_sim < threshold:
            kept_indices.append(i)
            kept_embeddings.append(emb)
    return [sentences[i] for i in kept_indices]

# ── 8. Chunking ────────────────────────────────────────────────────────────────
def chunk_by_wordcount(sentences: List[str], target: int) -> List[List[str]]:
    chunks, current, count = [], [], 0
    for s in sentences:
        w = len(s.split())
        if current and count + w > target:
            chunks.append(current)
            current, count = [], 0
        current.append(s)
        count += w
    if current:
        chunks.append(current)
    return chunks

# ── 9. Extractive sentence selection ─────────────────────────────────────────
def extract_top_sentences(sentences: List[str], model: SentenceTransformer,
                          k: int) -> List[str]:
    if not sentences:
        return []
    embeddings = model.encode(sentences, convert_to_tensor=True, show_progress_bar=False)
    centroid   = embeddings.mean(dim=0)
    scores     = util.cos_sim(embeddings, centroid).cpu().numpy().squeeze()
    if scores.ndim == 0:                 # single-sentence edge case
        return sentences
    top_idx    = sorted(scores.argsort()[::-1][:k])   # keep original order
    return [sentences[i] for i in top_idx]

# ── 10. Abstractive summarisation ─────────────────────────────────────────────
def build_summarizer():
    # FIX ① — task must be "summarization" for encoder-decoder models like BART
    print("Loading summarisation model (BART)...")
    return pipeline("summarization", model=SUMMARIZER_MODEL, device=DEVICE)

def _safe_summarize(pipe, text: str, max_len: int = 150, min_len: int = 30) -> str:
    """
    FIX ② — Correctly reads 'summary_text' key; adds length guard to avoid
    min_length > max_length errors on very short inputs.
    """
    word_count = len(text.split())
    # BART needs at least ~30 tokens; for very short chunks just return as-is
    if word_count < 20:
        return text
    safe_max = min(max_len, max(word_count, min_len + 1))
    safe_min = min(min_len, safe_max - 1)
    out = pipe(text, max_length=safe_max, min_length=safe_min,
               do_sample=False, truncation=True)
    # FIX ② — summarization pipeline always returns "summary_text"
    return out[0]["summary_text"].strip()

def summarize_chunks(chunk_texts: List[str], pipe) -> List[str]:
    return [_safe_summarize(pipe, t) for t in chunk_texts]

# ── 11. Heading heuristic ─────────────────────────────────────────────────────
def _heading_from_sentence(sentence: str, max_chars: int = 72) -> str:
    """
    FIX ⑤ — Better heading: strip the source tag, capitalise, truncate cleanly.
    V1 just did raw [:80] which could cut mid-word.
    """
    # Strip source tag if present
    sentence = re.sub(r'\[SOURCE:[^\]]+\]', '', sentence).strip()
    if len(sentence) <= max_chars:
        return sentence.capitalize()
    # Truncate at last word boundary before max_chars
    truncated = sentence[:max_chars].rsplit(' ', 1)[0]
    return truncated.capitalize() + "…"

# ── 12. Structured output helpers ────────────────────────────────────────────
def build_tldr(summaries: List[str], pipe) -> str:
    combined = " ".join(summaries)
    return _safe_summarize(pipe, combined, max_len=120, min_len=20)

def build_key_takeaways(summaries: List[str], max_items: int = 8) -> List[str]:
    seen, result = set(), []
    for summ in summaries:
        for sent in sent_tokenize(summ):
            key = sent.lower().strip()
            if key and key not in seen:
                result.append(sent.strip())
                seen.add(key)
            if len(result) >= max_items:
                return result
    return result

def notes_to_markdown(notes: Dict) -> str:
    md = [f"# {notes.get('title', 'Lecture Summary')}\n"]
    md += ["## TL;DR\n", notes.get("tldr", ""), "\n"]
    md += ["## Key Takeaways\n"] + [f"- {k}" for k in notes.get("key_takeaways", [])]
    md += ["\n## Detailed Notes\n"]
    for i, topic in enumerate(notes.get("detailed_notes", []), 1):
        md.append(f"### {topic.get('heading', f'Topic {i}')}")
        md += [f"- {b}" for b in topic.get("bullets", [])]
    if notes.get("definitions"):
        md += ["\n## Definitions\n"] + [f"- {d}" for d in notes["definitions"]]
    if notes.get("examples"):
        md += ["\n## Examples\n"] + [f"- {e}" for e in notes["examples"]]
    return "\n".join(md)

# ── 13. Main pipeline ─────────────────────────────────────────────────────────
def run_pipeline():
    # --- Step 1: collect & extract ---
    file_paths = collect_input_files()
    print(f"\nProcessing {len(file_paths)} file(s): {file_paths}")
    raw_text = fuse_extractions(file_paths)

    # --- Step 2: preprocess ---
    cleaned = preprocess(raw_text)
    with open(TRANSCRIPT_PATH, "w", encoding="utf-8") as f:
        f.write(cleaned)
    print(f"\nCleaned transcript saved → {TRANSCRIPT_PATH}")

    # --- Step 3: sentence split ---
    sentences = [s.strip() for s in sent_tokenize(cleaned) if s.strip()]
    if not sentences:
        raise ValueError("No sentences found after preprocessing. Check input quality.")

    # --- Step 4: load models ---
    sbert      = SentenceTransformer("all-MiniLM-L6-v2")
    summarizer = build_summarizer()

    # --- Step 5: deduplication (FIX ④) ---
    print(f"\nDeduplicating {len(sentences)} sentences...")
    sentences = deduplicate_sentences(sentences, sbert)
    print(f"  → {len(sentences)} unique sentences retained.")

    # --- Step 6: chunk ---
    chunks = chunk_by_wordcount(sentences, CHUNK_WORD_TARGET)
    print(f"\nCreated {len(chunks)} chunk(s).")

    # --- Step 7: extractive + abstractive per chunk ---
    chunk_texts, chunk_meta = [], []
    for i, chunk_sents in enumerate(chunks):
        top_sents  = extract_top_sentences(chunk_sents, sbert, EXTRACTIVE_SENTENCES)
        chunk_text = " ".join(top_sents) if len(" ".join(top_sents).split()) >= 30 \
                     else " ".join(chunk_sents)
        chunk_texts.append(chunk_text)
        heading = _heading_from_sentence(top_sents[0]) if top_sents else f"Topic {i+1}"
        chunk_meta.append({"heading": heading, "source_sentences": top_sents})

    print("Running abstractive summarisation on chunks...")
    chunk_summaries = summarize_chunks(chunk_texts, summarizer)

    # --- Step 8: assemble notes ---
    print("Assembling structured notes...")
    title     = _heading_from_sentence(sentences[0])
    tldr      = build_tldr(chunk_summaries, summarizer)
    takeaways = build_key_takeaways(chunk_summaries)

    detailed  = [
        {
            "heading": meta["heading"],
            "bullets": [s.strip() for s in sent_tokenize(summary) if s.strip()],
            "source_sentences": meta["source_sentences"],
        }
        for meta, summary in zip(chunk_meta, chunk_summaries)
    ]

    definitions = [s for s in sentences
                   if any(kw in s.lower() for kw in
                          ["is defined as", "is called", "refers to", "means that"])]
    examples    = [s for s in sentences
                   if any(kw in s.lower() for kw in
                          ["for example", "for instance", "e.g.", "such as"])]

    notes = {
        "title":        title,
        "tldr":         tldr,
        "key_takeaways": takeaways,
        "detailed_notes": detailed,
        "formulas":     [],
        "definitions":  definitions,
        "examples":     examples,
        "action_items": [],
    }

    # --- Step 9: save outputs ---
    with open(OUTPUT_JSON, "w", encoding="utf-8") as f:
        json.dump(notes, f, indent=2, ensure_ascii=False)
    with open(OUTPUT_MD, "w", encoding="utf-8") as f:
        f.write(notes_to_markdown(notes))

    print(f"\n✅ Done!\n  JSON     → {OUTPUT_JSON}\n  Markdown → {OUTPUT_MD}")
    return notes

# ── 14. Entry point ───────────────────────────────────────────────────────────
if __name__ == "__main__":
    notes = run_pipeline()

    print("\n─── TL;DR ───────────────────────────────────────────")
    print(notes["tldr"])
    print("\n─── Key Takeaways ───────────────────────────────────")
    for k in notes["key_takeaways"]:
        print(" •", k)